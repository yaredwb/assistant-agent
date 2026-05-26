"""PowerPoint companion deck generation for review walkthroughs.

The deck is intentionally built from the existing dossier contract. It is a visual
agenda for the voice walkthrough, not a separate source of truth.
"""

from __future__ import annotations

import json
import shutil
from dataclasses import asdict, dataclass
from datetime import datetime, timezone
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt

from .contract import ChangedFile, Dossier
from .dossier import reviews_dir


MAX_BULLETS = 7


@dataclass(frozen=True)
class SlideSpec:
    title: str
    bullets: list[str]
    narration: str


@dataclass(frozen=True)
class DeckResult:
    pptx_path: Path
    outline_path: Path
    latest_pptx_path: Path
    latest_outline_path: Path
    slide_count: int


def _now_stamp() -> str:
    return datetime.now(timezone.utc).strftime("%Y%m%dT%H%M%SZ")


def _clean(text: str | None, *, fallback: str = "") -> str:
    text = (text or "").strip()
    return " ".join(text.split()) if text else fallback


def _truncate(text: str, limit: int = 170) -> str:
    text = _clean(text)
    if len(text) <= limit:
        return text
    return text[: limit - 3].rstrip() + "..."


def _limited(items: list[str], *, fallback: str, limit: int = MAX_BULLETS) -> list[str]:
    cleaned = [_truncate(item) for item in items if _clean(item)]
    if not cleaned:
        return [fallback]
    if len(cleaned) <= limit:
        return cleaned
    remaining = len(cleaned) - limit
    return cleaned[:limit] + [f"Plus {remaining} more."]


def _file_bullet(changed_file: ChangedFile) -> str:
    stats: list[str] = []
    if changed_file.additions is not None:
        stats.append(f"+{changed_file.additions}")
    if changed_file.deletions is not None:
        stats.append(f"-{changed_file.deletions}")
    stat_text = f" ({', '.join(stats)})" if stats else ""
    summary = f": {_truncate(changed_file.summary, 110)}" if changed_file.summary else ""
    return f"{changed_file.path}{stat_text}{summary}"


def build_slide_outline(dossier: Dossier) -> list[SlideSpec]:
    """Convert a dossier into a concise slide outline with narration cues."""
    generated = _clean(dossier.generated_at, fallback="not recorded")
    branch = _clean(dossier.branch, fallback="not recorded")
    task = _clean(dossier.task, fallback="Task not recorded in the dossier.")
    summary = _clean(dossier.summary, fallback="Summary not recorded in the dossier.")
    changed_files = [_file_bullet(item) for item in dossier.changed_files]
    commits = dossier.evidence.commits
    diff_refs = dossier.evidence.diff_refs
    transcript_refs = dossier.evidence.transcript_refs

    slides = [
        SlideSpec(
            title=f"{dossier.repo} review",
            bullets=[
                f"Branch: {branch}",
                f"Task: {_truncate(task, 210)}",
                f"Dossier generated: {generated}",
            ],
            narration="Open by framing the repo, branch, and task that was reviewed.",
        ),
        SlideSpec(
            title="Walkthrough agenda",
            bullets=_limited(
                dossier.suggested_walkthrough_order,
                fallback="Use the standard flow: summary, changes, decisions, tests, risks, and questions.",
            ),
            narration="Use this slide as the agenda and pause for questions after each major topic.",
        ),
        SlideSpec(
            title="Executive summary",
            bullets=[_truncate(summary, 360)],
            narration="Give the one-paragraph headline before moving into implementation detail.",
        ),
        SlideSpec(
            title="What changed",
            bullets=_limited(changed_files, fallback="No changed files were listed in the dossier."),
            narration="Summarize the main changed files and explain how they fit together.",
        ),
        SlideSpec(
            title="Key decisions",
            bullets=_limited(
                dossier.key_decisions,
                fallback="No explicit implementation decisions were listed.",
            ),
            narration="Explain why the work was shaped this way, not just what files changed.",
        ),
        SlideSpec(
            title="Tests and verification",
            bullets=_limited(
                dossier.tests_run,
                fallback="No tests or verification steps were recorded.",
            ),
            narration="Be precise about what was verified and what was not verified.",
        ),
        SlideSpec(
            title="Risks to review",
            bullets=_limited(dossier.risks, fallback="No risks were recorded in the dossier."),
            narration=(
                "Call out the parts that deserve the engineer's attention before they "
                "share or build on the work."
            ),
        ),
        SlideSpec(
            title="Open questions",
            bullets=_limited(
                dossier.open_questions,
                fallback="No open questions were recorded.",
            ),
            narration="Use these questions to invite feedback and decide whether follow-up work is needed.",
        ),
        SlideSpec(
            title="Evidence pointers",
            bullets=_limited(
                [
                    *(f"Commit: {commit}" for commit in commits),
                    *(f"Diff: {ref}" for ref in diff_refs),
                    *(f"Transcript: {ref}" for ref in transcript_refs),
                ],
                fallback="No commit, diff, or transcript pointers were recorded.",
            ),
            narration="Use these pointers when the engineer asks for proof or deeper context.",
        ),
        SlideSpec(
            title="Feedback and follow-up",
            bullets=[
                "Capture requested changes during the conversation.",
                "Turn feedback into a follow-up prompt for the coding agent.",
                "Keep the voice agent focused on explanation, Q&A, and handoff.",
            ],
            narration="Close by asking what should change, then record concrete feedback items.",
        ),
    ]
    return slides


def outline_to_json(slides: list[SlideSpec]) -> str:
    payload = {
        "generated_at": datetime.now(timezone.utc).isoformat(timespec="seconds"),
        "slides": [
            {"number": index, **asdict(slide)}
            for index, slide in enumerate(slides, start=1)
        ],
    }
    return json.dumps(payload, indent=2)


def _add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: int,
    bold: bool = False,
    color=(30, 38, 48),
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return box


def _add_bullets(slide, bullets: list[str], *, top=Inches(1.75)):
    # build_slide_outline is the single truncation authority; render bullets as sized
    # (don't re-truncate here, or the per-slide limits — e.g. the 360-char executive
    # summary — get silently clipped).
    box = slide.shapes.add_textbox(Inches(0.82), top, Inches(11.35), Inches(4.7))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.05)
    frame.margin_right = Inches(0.05)
    frame.margin_top = Inches(0.05)
    frame.margin_bottom = Inches(0.05)

    font_size = 20 if len(bullets) <= 4 else 17
    for index, item in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"- {item}"
        paragraph.level = 0
        paragraph.space_after = Pt(9)
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = RGBColor(38, 46, 56)
    return box


def _add_footer(slide, number: int, total: int):
    _add_textbox(
        slide,
        Inches(0.72),
        Inches(6.92),
        Inches(7.5),
        Inches(0.28),
        "AI work review walkthrough",
        size=9,
        color=(105, 115, 128),
    )
    _add_textbox(
        slide,
        Inches(11.55),
        Inches(6.92),
        Inches(0.95),
        Inches(0.28),
        f"{number}/{total}",
        size=9,
        color=(105, 115, 128),
    )


def render_pptx(slides: list[SlideSpec], output_path: Path) -> Path:
    """Render slide specs to a PowerPoint file."""
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]
    total = len(slides)
    for number, spec in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank)
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = RGBColor(248, 250, 252)

        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(29, 78, 216)
        accent.line.fill.background()

        _add_textbox(
            slide,
            Inches(0.72),
            Inches(0.42),
            Inches(11.8),
            Inches(0.7),
            spec.title,
            size=30,
            bold=True,
            color=(15, 23, 42),
        )
        _add_bullets(slide, spec.bullets)
        _add_footer(slide, number, total)

    prs.save(output_path)
    return output_path


def write_deck(repo_path: Path, dossier: Dossier) -> DeckResult:
    """Create timestamped and latest PPTX/outline artifacts for a repo review."""
    slides = build_slide_outline(dossier)
    review_dir = reviews_dir(repo_path)
    stamp = _now_stamp()

    pptx_path = review_dir / f"deck-{stamp}.pptx"
    outline_path = review_dir / f"deck-{stamp}.outline.json"
    latest_pptx_path = review_dir / "latest.pptx"
    latest_outline_path = review_dir / "latest.deck.json"

    render_pptx(slides, pptx_path)
    outline_path.write_text(outline_to_json(slides), encoding="utf-8")

    shutil.copyfile(pptx_path, latest_pptx_path)
    shutil.copyfile(outline_path, latest_outline_path)

    return DeckResult(
        pptx_path=pptx_path,
        outline_path=outline_path,
        latest_pptx_path=latest_pptx_path,
        latest_outline_path=latest_outline_path,
        slide_count=len(slides),
    )
