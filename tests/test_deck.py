import json

from pptx import Presentation

from assistant_agent.contract import ChangedFile, Dossier, Evidence
from assistant_agent.deck import build_slide_outline, write_deck


def _dossier() -> Dossier:
    return Dossier(
        repo="demo",
        branch="main",
        task="Add deck generation",
        summary="Created a PowerPoint companion deck for voice walkthroughs.",
        changed_files=[
            ChangedFile(path="src/assistant_agent/deck.py", summary="Build and render slides."),
            ChangedFile(path="src/assistant_agent/cli.py", summary="Wire deck command into the CLI."),
        ],
        key_decisions=["Keep the deck as a visual agenda, not a source of truth."],
        tests_run=["uv run pytest -q"],
        risks=["Slide content can become stale if the dossier is stale."],
        open_questions=["Should the browser UI show the deck during voice calls?"],
        suggested_walkthrough_order=["deck", "voice", "feedback"],
        evidence=Evidence(commits=["abc123"], diff_refs=["file:src/assistant_agent/deck.py"]),
    )


def test_build_slide_outline_has_walkthrough_structure():
    outline = build_slide_outline(_dossier())

    assert outline[0].title == "demo review"
    assert any(slide.title == "Walkthrough agenda" for slide in outline)
    assert any(slide.title == "Feedback and follow-up" for slide in outline)
    assert any("slide" in slide.narration.lower() for slide in outline)


def test_write_deck_creates_pptx_and_outline(tmp_path):
    result = write_deck(tmp_path, _dossier())

    assert result.pptx_path.exists()
    assert result.latest_pptx_path.exists()
    assert result.outline_path.exists()
    assert result.latest_outline_path.exists()

    deck = Presentation(result.pptx_path)
    assert len(deck.slides) == result.slide_count
    assert result.slide_count >= 8

    outline = json.loads(result.outline_path.read_text())
    assert outline["slides"][0]["title"] == "demo review"
    assert outline["slides"][0]["number"] == 1


def _slide_body(slide) -> str:
    """The longest text frame on a slide is its bullet body."""
    texts = [sh.text_frame.text for sh in slide.shapes if sh.has_text_frame]
    return max(texts, key=len) if texts else ""


def test_executive_summary_is_not_over_truncated(tmp_path):
    # Regression: _add_bullets used to re-truncate every bullet to 155 chars, clipping
    # the 360-char executive-summary slide. build_slide_outline is the sole authority now.
    long_summary = "Sentence number {}. ".format
    dossier = _dossier()
    dossier.summary = "".join(long_summary(i) for i in range(1, 60))  # ~600+ chars
    assert len(dossier.summary) > 360

    result = write_deck(tmp_path, dossier)
    deck = Presentation(result.pptx_path)
    exec_slide = next(
        s for s in deck.slides
        if any(sh.has_text_frame and sh.text_frame.text == "Executive summary" for sh in s.shapes)
    )
    body = _slide_body(exec_slide)
    # Should reflect the 360-char per-slide limit, not the old 155 cap.
    assert len(body) > 300
